from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "teslim"
OUT.mkdir(exist_ok=True)

PPTX_PATH = OUT / "sunum_3_kisilik_xai.pptx"

NAVY = RGBColor(31, 78, 121)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(15, 118, 110)
ORANGE = RGBColor(249, 115, 22)
PURPLE = RGBColor(109, 40, 217)
RED = RGBColor(190, 18, 60)
GRAY = RGBColor(75, 85, 99)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(17, 24, 39)


def clear_frame(tf):
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)


def add_text(slide, x, y, w, h, text, size=16, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    clear_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    return shape


def add_paragraphs(slide, x, y, w, h, paragraphs, size=14.2, color=BLACK, spacing=5):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    clear_frame(tf)
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
    return shape


def add_bullets(slide, x, y, w, h, bullets, size=13.7, color=BLACK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    clear_frame(tf)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4.5)
        p.level = 0
    return shape


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.55, 0.22, 12.2, 0.50, title, 22.5, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.86), Inches(12.2), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if subtitle:
        add_text(slide, 0.58, 0.94, 11.8, 0.38, subtitle, 11.8, GRAY)


def add_footer(slide, n):
    add_text(slide, 0.55, 7.18, 8.8, 0.22, "BLM308 Veri Madenciliği Final Projesi", 8.3, GRAY)
    add_text(slide, 12.2, 7.18, 0.5, 0.22, str(n), 8.3, GRAY, align=PP_ALIGN.RIGHT)


def content_slide(prs, title, n, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    add_footer(slide, n)
    return slide


def section_slide(prs, title, subtitle, paragraphs, n, color=NAVY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color
    add_text(slide, 0.75, 1.45, 11.8, 0.70, title, 34, WHITE, True)
    add_text(slide, 0.78, 2.32, 11.4, 0.60, subtitle, 18.5, RGBColor(219, 234, 254))
    add_paragraphs(slide, 0.80, 3.23, 11.05, 1.72, paragraphs, 14.6, WHITE, 6)
    add_footer(slide, n)
    return slide


def add_callout(slide, x, y, w, h, title, body, color=BLUE, title_size=13.4, body_size=11.8):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = color
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    clear_frame(tf)
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(body_size)
    p2.font.color.rgb = BLACK
    p2.space_before = Pt(4)
    return shape


def add_table(slide, x, y, w, h, data, font_size=10.4):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for r_idx, row in enumerate(data):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r_idx == 0 else RGBColor(248, 250, 252)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE if r_idx == 0 else BLACK
                p.alignment = PP_ALIGN.CENTER
    return shape


def add_image(slide, image_name, x, y, w=None, h=None):
    path = FIG / image_name
    if path.exists():
        kwargs = {}
        if w:
            kwargs["width"] = Inches(w)
        if h:
            kwargs["height"] = Inches(h)
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)
    return add_callout(slide, x, y, w or 5, h or 2, "Görsel bulunamadı", image_name, RED)


def add_step_bar(slide, labels, y=2.24):
    colors = [NAVY, BLUE, TEAL, ORANGE, PURPLE, RGBColor(14, 116, 144)]
    for i, label in enumerate(labels):
        x = 0.58 + i * 2.08
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.62), Inches(0.86))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        shape.text_frame.clear()
        add_text(slide, x + 0.04, y + 0.13, 1.54, 0.52, label, 12.3, WHITE, True, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            add_text(slide, x + 1.62, y + 0.25, 0.34, 0.24, "→", 17, GRAY, True, PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(239, 246, 255)
    add_text(slide, 0.70, 0.62, 11.8, 1.45, "Banka Pazarlama Kampanyalarında\nMüşteri Tahmin ve Açıklanabilirlik Sistemi", 33, NAVY, True)
    add_text(slide, 0.78, 2.28, 9.8, 0.42, "BLM308 Veri Madenciliği - Final Proje", 18, GRAY)
    add_paragraphs(slide, 0.78, 3.02, 10.9, 1.45, [
        "Bu çalışma, banka pazarlama kampanyalarında hangi müşterilerin vadeli mevduata abone olma olasılığının daha yüksek olduğunu tahmin eden ve bu tahminleri açıklanabilir hale getiren uçtan uca bir veri madenciliği projesidir.",
        "Proje yalnızca model başarısına odaklanmaz; veri sızıntısı riski, sınıf dengesizliği, istatistiksel karşılaştırma, SHAP/LIME temelli açıklanabilirlik ve Türkçe uygulama arayüzü birlikte ele alınır."
    ], 14.2)
    add_text(slide, 0.78, 5.80, 10.8, 0.34, "Sude ANLAŞ · Doğan Fatih OĞUR · Merve Naz ORAN", 16.5, BLACK, True)
    add_text(slide, 0.78, 6.20, 8.6, 0.28, "Gedik Üniversitesi · Bilgisayar Mühendisliği · Bahar 2026", 12.3, GRAY)
    for i, (label, col) in enumerate([("Veri", BLUE), ("Model", TEAL), ("Açıklama", ORANGE)]):
        x = 8.20 + i * 1.48
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.45), Inches(1.06), Inches(1.06))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        add_text(slide, x - 0.05, 5.62, 1.18, 0.28, label, 11.3, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, n)
    n += 1

    slide = content_slide(prs, "Sunum Akışı", n, "Anlatım, iş probleminden başlayıp model çıktılarının gerçek kullanımda nasıl yorumlanacağına kadar ilerler.")
    add_table(slide, 0.65, 1.35, 12.0, 4.2, [
        ["Bölüm", "Bu bölümde cevaplanan temel soru"],
        ["1. Problem ve motivasyon", "Banka neden veriye dayalı müşteri önceliklendirmesine ihtiyaç duyar?"],
        ["2. Yöntem", "CRISP-DM süreci projede hangi adımlarla uygulanmıştır?"],
        ["3. Veri", "Veri seti hangi riskleri taşır ve nasıl hazırlanmıştır?"],
        ["4. Modelleme", "Hangi modeller denenmiş, ana model neden seçilmiştir?"],
        ["5. Sonuçlar", "Başarı metrikleri ve istatistiksel testler ne göstermektedir?"],
        ["6. Açıklanabilirlik", "Model kararları genel ve tek müşteri düzeyinde nasıl açıklanmıştır?"],
        ["7. Uygulama", "Çalışma nasıl gösterilebilir bir Türkçe karar destek sistemine dönüştürülmüştür?"],
    ], 10.5)
    add_paragraphs(slide, 0.95, 5.90, 11.5, 0.70, [
        "Bu yapı, projenin yalnızca teknik bir model denemesi olmadığını; veri madenciliği döngüsünün tüm ana bileşenlerini kapsayan bütüncül bir çalışma olduğunu gösterir."
    ], 14.0, GRAY)
    n += 1

    section_slide(prs, "1. Problem ve Motivasyon", "Tahmin sistemi, kampanya kaynaklarını daha etkili kullanmak ve model kararlarını denetlenebilir hale getirmek için tasarlanmıştır.", [
        "Bankacılıkta her müşteriye aynı sıklıkta ve aynı öncelikle ulaşmak maliyetli bir yaklaşımdır. Bu nedenle kampanya ekipleri, sınırlı zaman ve arama kapasitesini dönüşüm ihtimali daha yüksek müşterilere yönlendirmek ister.",
        "Ancak finansal alanda yalnızca yüksek tahmin skoru yeterli değildir; modelin bu skoru hangi müşteri özelliklerine dayanarak verdiği de açıklanmalıdır."
    ], n)
    n += 1

    slide = content_slide(prs, "Problem Tanımı", n)
    add_paragraphs(slide, 0.75, 1.32, 6.15, 3.45, [
        "Projenin temel problemi, bir müşterinin banka pazarlama kampanyası sonucunda vadeli mevduata abone olup olmayacağını önceden tahmin etmektir. Bu tahmin, kampanya başlamadan önce müşteri önceliklendirmesi yapmak için kullanılır.",
        "Rastgele arama yaklaşımında banka, abonelik potansiyeli düşük müşterilere de aynı kaynakları ayırır. Bu durum hem operasyonel maliyeti artırır hem de kampanya ekibinin verimini düşürür.",
        "Modelin amacı, müşterilere ait profil, finansal durum ve kampanya geçmişi bilgilerinden yararlanarak yüksek olasılıklı müşterileri öne çıkarmaktır."
    ], 14.0)
    add_callout(slide, 7.35, 1.45, 4.8, 1.35, "Karar destek sorusu", "“Hangi müşteriler önce aranmalı ve bu öneri hangi değişkenlere dayanmaktadır?”", BLUE, 14.0, 12.8)
    add_callout(slide, 7.35, 3.18, 4.8, 1.72, "Gerçek kullanım senaryosu", "Bir kampanya yöneticisi, günlük arama listesinde en yüksek olasılıklı müşterileri üst sıraya alır; XAI çıktıları ise bu sıralamanın neden oluştuğunu açıklar.", TEAL, 14.0, 12.0)
    n += 1

    slide = content_slide(prs, "İş Değeri ve Beklenen Katkı", n)
    add_table(slide, 0.65, 1.25, 12.0, 2.35, [
        ["Klasik yaklaşım", "Veriye dayalı yaklaşım"],
        ["Müşteriler çoğunlukla genel liste sırasına göre aranır.", "Müşteriler abonelik olasılığına göre önceliklendirilir."],
        ["Arama başına dönüşüm oranı düşük kalabilir.", "Kampanya ekibi daha verimli temas kurabilir."],
        ["Kararın gerekçesi çoğu zaman sezgiseldir.", "Tahmin sonucu açıklanabilir faktörlerle desteklenir."],
    ], 11.3)
    add_paragraphs(slide, 0.85, 4.05, 11.65, 1.55, [
        "Bu katkı yalnızca satış performansıyla sınırlı değildir. Model, kampanya planlamasında veri temelli karar alma kültürünü destekler; hangi müşteri segmentlerinin daha duyarlı olduğunu görünür hale getirir.",
        "Açıklanabilirlik katmanı, pazarlama uzmanının model önerisini körü körüne kabul etmesini değil, iş bilgisiyle birlikte yorumlamasını sağlar."
    ], 14.2)
    n += 1

    slide = content_slide(prs, "Neden Açıklanabilir Yapay Zeka?", n)
    add_paragraphs(slide, 0.78, 1.25, 6.2, 3.85, [
        "Random Forest veya çok katmanlı algılayıcı gibi modeller karmaşık ilişkileri öğrenebilir; ancak karar mantığı doğrudan okunamadığı için kara kutu olarak değerlendirilir.",
        "Bankacılık gibi denetlenebilirliğin önemli olduğu alanlarda modelin yalnızca “yüksek olasılık” demesi yeterli değildir. Kararın hangi değişkenlerle desteklendiği, hangi faktörlerin sonucu düşürdüğü veya artırdığı açıklanmalıdır.",
        "Bu projede XAI, ek bir görselleştirme adımı değil; modelin güvenilir, anlaşılır ve sunulabilir olmasını sağlayan temel bileşendir."
    ], 13.8)
    add_callout(slide, 7.35, 1.45, 4.75, 1.48, "Küresel açıklama", "Modelin genel olarak hangi değişkenlere daha fazla önem verdiğini gösterir. Örneğin ekonomik göstergeler veya önceki kampanya sonucu öne çıkabilir.", ORANGE, 13.6, 11.7)
    add_callout(slide, 7.35, 3.25, 4.75, 1.62, "Yerel açıklama", "Tek bir müşteri için tahmini hangi özelliklerin artırdığını veya azalttığını gösterir. Bu, müşteri bazlı karar gerekçesi üretir.", PURPLE, 13.6, 11.7)
    n += 1

    slide = content_slide(prs, "Proje Kapsamı ve Güçlü Yönler", n)
    add_bullets(slide, 0.85, 1.25, 11.8, 4.65, [
        "Çalışma, veri temizleme ve model eğitimiyle sınırlı bırakılmamış; rapor, sunum, kod yapısı ve çalışan Streamlit arayüzüyle teslim edilebilir bir proje bütününe dönüştürülmüştür.",
        "Görüşme süresi değişkeni için veri sızıntısı riski açıkça tartışılmıştır. Bu karar, modelin gerçek kampanya öncesi senaryoda kullanılabilirliğini güçlendirir.",
        "Random Forest ve çok katmanlı algılayıcı gibi kara kutu modeller, açıklanabilirlik araçlarıyla desteklenmiştir. Böylece proje, seçilen XAI başlığıyla doğrudan uyumludur.",
        "McNemar testi eklenerek model karşılaştırması yalnızca yüzeysel metriklerle değil, istatistiksel anlamlılık açısından da değerlendirilmiştir.",
        "Türkçe arayüz, hocanın raporda anlatılan sistemi uygulamalı olarak görmesine imkân verir."
    ], 14.1)
    n += 1

    section_slide(prs, "2. Yöntem", "Proje, CRISP-DM yaklaşımına göre yapılandırılmış ve her aşama bir sonraki kararı besleyecek şekilde ilerletilmiştir.", [
        "CRISP-DM yöntemi, veri madenciliği projelerinde iş anlayışı, veri anlayışı, veri hazırlığı, modelleme, değerlendirme ve yaygınlaştırma adımlarını sistematik biçimde ele alır.",
        "Bu projede metodoloji yalnızca teorik olarak yazılmamış; veri setinin incelenmesinden Türkçe uygulama arayüzüne kadar tüm adımlarda uygulanmıştır."
    ], n, BLUE)
    n += 1

    slide = content_slide(prs, "Genel Süreç: CRISP-DM Uygulaması", n)
    add_step_bar(slide, ["İş\nAnlayışı", "Veri\nAnlayışı", "Veri\nHazırlığı", "Model\nEğitimi", "Değer-\nlendirme", "XAI ve\nArayüz"])
    add_table(slide, 0.75, 3.65, 11.85, 2.05, [
        ["Aşama", "Projede karşılığı"],
        ["İş anlayışı", "Kampanya önceliklendirme ve müşteri tahmini problemi tanımlandı."],
        ["Veri hazırlığı", "Kategorik değişkenler dönüştürüldü, sayısal değişkenler standartlaştırıldı."],
        ["Modelleme", "Temel ve kara kutu modeller aynı test yapısında karşılaştırıldı."],
        ["Değerlendirme", "Metrikler, hata matrisi, ROC-AUC ve McNemar testi birlikte yorumlandı."],
    ], 10.9)
    add_paragraphs(slide, 0.90, 6.05, 11.5, 0.45, [
        "Bu süreç, model seçiminin rastgele değil; veri özellikleri, iş ihtiyacı ve açıklanabilirlik gereksinimiyle bağlantılı yapılmasını sağlar."
    ], 13.4, GRAY)
    n += 1

    slide = content_slide(prs, "Kullanılan Teknolojiler ve Proje Ortamı", n)
    add_table(slide, 0.65, 1.25, 12.0, 2.75, [
        ["Bileşen", "Kullanım amacı"],
        ["Python, pandas, numpy", "Veri okuma, dönüştürme ve sayısal işlemler"],
        ["scikit-learn", "Modelleme, ön işleme, metrik hesaplama ve test bölme"],
        ["matplotlib, seaborn", "Grafik ve performans çıktılarının üretilmesi"],
        ["SHAP/LIME yaklaşımı", "Küresel ve yerel model açıklamalarının oluşturulması"],
        ["Streamlit", "Hocaya gösterilebilir Türkçe tahmin arayüzü"],
    ], 10.9)
    add_paragraphs(slide, 0.85, 4.45, 11.65, 1.25, [
        "Proje Visual Studio Code + Python ortamında çalıştırılacak şekilde düzenlenmiştir. Bu tercih, kodun okunmasını, dosya yapısının incelenmesini ve uygulamanın doğrudan başlatılmasını kolaylaştırır.",
        "Ayrıca train.py, explain.py ve app.py dosyaları ayrı tutulduğu için eğitim, açıklama üretimi ve arayüz çalıştırma adımları net biçimde ayrılmıştır."
    ], 13.7)
    n += 1

    section_slide(prs, "3. Veri Seti ve Hazırlık", "Veri hazırlığı aşaması, model performansını doğrudan etkilediği için proje içinde kritik bir karar alanı olarak ele alınmıştır.", [
        "UCI Bank Marketing veri seti, pazarlama kampanyası sonucunda müşterinin vadeli mevduata abone olup olmadığını gösteren gerçekçi bir sınıflandırma problemidir.",
        "Bu veri setinde sınıf dengesizliği ve görüşme süresi değişkeni gibi dikkat edilmesi gereken noktalar vardır. Bu nedenle veri hazırlığı yalnızca teknik dönüşüm değil, aynı zamanda metodolojik bir güvenilirlik aşamasıdır."
    ], n, TEAL)
    n += 1

    slide = content_slide(prs, "Veri Kaynağı ve Temel Yapı", n)
    add_table(slide, 0.70, 1.25, 5.75, 2.45, [
        ["Özellik", "Değer"],
        ["Kaynak", "UCI Bank Marketing"],
        ["Örnek sayısı", "41.188"],
        ["Girdi değişkeni", "20"],
        ["Hedef değişken", "Vadeli mevduat aboneliği"],
        ["Pozitif sınıf", "Abone oldu"],
    ], 11.2)
    add_paragraphs(slide, 6.95, 1.28, 5.35, 2.72, [
        "Veri seti, müşteri profili, finansal durum, kampanya temas bilgileri ve ekonomik göstergelerden oluşur. Bu çeşitlilik, modelin yalnızca tek bir müşteri özelliğine değil, birden fazla karar sinyaline bakmasını sağlar.",
        "Hedef değişken ikili yapıdadır: müşteri kampanya sonunda vadeli mevduata abone olmuştur veya olmamıştır. Bu nedenle problem, denetimli ikili sınıflandırma problemi olarak modellenmiştir."
    ], 13.1)
    add_callout(slide, 0.90, 4.60, 11.5, 1.05, "Akademik değerlendirme", "Veri kaynağının açık olması, çalışmanın tekrar üretilebilirliğini artırır. Aynı veri seti üzerinde farklı modeller denenebilir, sonuçlar karşılaştırılabilir ve model kararları daha şeffaf biçimde tartışılabilir.", BLUE, 13.6, 12.2)
    n += 1

    slide = content_slide(prs, "Hedef Değişken Dağılımı", n)
    add_image(slide, "target_distribution_turkce.png", 1.03, 1.17, w=5.85)
    add_paragraphs(slide, 7.20, 1.30, 5.15, 3.90, [
        "Grafik, abone olmayan müşterilerin veri setinde açık biçimde çoğunlukta olduğunu gösterir. Bu durum sınıf dengesizliği problemidir ve model değerlendirmesinde doğrudan dikkate alınmalıdır.",
        "Eğer yalnızca doğruluk metriğine bakılırsa model çoğunluk sınıfını tahmin ederek yüksek görünen fakat iş açısından zayıf bir başarı elde edebilir. Bu nedenle kesinlik, duyarlılık, F1 ve ROC-AUC metrikleri birlikte raporlanmıştır.",
        "Bankanın hedefi yalnızca doğru tahmin sayısını artırmak değil, abone olabilecek müşterileri yakalayabilmektir. Bu yüzden pozitif sınıf performansı özellikle önemlidir."
    ], 13.2)
    n += 1

    slide = content_slide(prs, "Veri Hazırlığı Kararları", n)
    add_table(slide, 0.65, 1.20, 12.0, 3.08, [
        ["Karar", "Neden gerekli?"],
        ["Kategorik değişkenlerin dönüştürülmesi", "Meslek, medeni durum ve iletişim kanalı gibi alanlar sayısal modele uygun hale getirildi."],
        ["Sayısal değişkenlerin ölçeklenmesi", "Özellikle sinir ağı modelinde değişken ölçeklerinin öğrenmeyi bozması engellendi."],
        ["Stratified train-test split", "Azınlık sınıf oranının eğitim ve test kümelerinde korunması sağlandı."],
        ["Sınıf ağırlığı kullanımı", "Abone olan azınlık sınıfın model tarafından tamamen göz ardı edilmesi azaltıldı."],
    ], 10.7)
    add_paragraphs(slide, 0.85, 4.73, 11.65, 1.05, [
        "Bu kararların ortak amacı, modelin yalnızca eğitim verisine ezberci biçimde uyum sağlamasını değil, daha önce görmediği test örnekleri üzerinde de tutarlı davranmasını sağlamaktır.",
        "Veri hazırlığı aşaması bu nedenle projenin teknik temelini oluşturur; hatalı hazırlanan veri, daha sonra kullanılan gelişmiş modelin de yanlış sonuç üretmesine yol açabilir."
    ], 13.3)
    n += 1

    slide = content_slide(prs, "Kritik Nokta: Görüşme Süresi ve Veri Sızıntısı", n)
    add_paragraphs(slide, 0.78, 1.25, 6.25, 4.05, [
        "Görüşme süresi değişkeni, müşteriyle yapılan aramanın ne kadar sürdüğünü ifade eder. Bu değişken genellikle hedef değişkenle güçlü ilişki taşır; çünkü uzun görüşmeler çoğu zaman müşterinin daha ilgili olduğunu gösterebilir.",
        "Fakat gerçek kampanya öncesi tahmin senaryosunda bu bilgi henüz mevcut değildir. Banka müşteriyi aramadan önce görüşmenin kaç dakika süreceğini bilemez.",
        "Bu nedenle görüşme süresini kullanmak, modelin testte yüksek performans göstermesine rağmen gerçek kullanımda uygulanamayacak bir bilgiye dayanması anlamına gelir. Projede bu risk açıkça belirtilmiş ve gerçekçi senaryoda değişken modelden çıkarılmıştır."
    ], 13.5)
    add_callout(slide, 7.50, 1.45, 4.55, 1.70, "Neden önemli?", "Veri sızıntısı, modelin gelecekte bilinmeyecek bilgileri öğrenmesine neden olur. Bu durum akademik olarak yanıltıcı ve uygulama açısından güvenilmez sonuç üretir.", RED, 13.8, 11.8)
    add_callout(slide, 7.50, 3.55, 4.55, 1.45, "Projede alınan karar", "Hem görüşme süresi dahil hem de hariç senaryolar raporlandı; ana yorum gerçekçi senaryo olan görüşme süresi hariç modele dayandırıldı.", TEAL, 13.8, 11.8)
    n += 1

    section_slide(prs, "4. Modelleme Yaklaşımı", "Modelleme aşamasında basit referans modeller ile kara kutu modeller birlikte değerlendirilmiştir.", [
        "Bu yaklaşım, yalnızca en karmaşık modeli seçmek yerine model başarısını karşılaştırmalı olarak görmeyi sağlar. Basit modeller referans noktası oluştururken, kara kutu modeller daha karmaşık ilişkileri yakalayabilir.",
        "Proje başlığı XAI olduğu için kara kutu modellerin açıklanabilir hale getirilmesi özellikle önemlidir."
    ], n, ORANGE)
    n += 1

    slide = content_slide(prs, "Seçilen Modeller ve Gerekçeleri", n)
    add_table(slide, 0.55, 1.20, 12.25, 3.35, [
        ["Model", "Projedeki rolü", "Gerekçe"],
        ["Lojistik Regresyon", "Referans model", "Basit, hızlı ve temel performans çizgisi sağlar."],
        ["Karar Ağacı", "Yorumlanabilir karşılaştırma", "Kural benzeri yapı ile ağaç tabanlı öğrenmeyi temsil eder."],
        ["Random Forest", "Ana kara kutu model", "Birden fazla ağacı birleştirerek daha kararlı tahmin üretir."],
        ["Çok Katmanlı Algılayıcı", "Alternatif kara kutu model", "Doğrusal olmayan ilişkileri öğrenebilen sinir ağı yaklaşımıdır."],
    ], 10.4)
    add_paragraphs(slide, 0.85, 5.02, 11.6, 0.92, [
        "Random Forest, performans ve açıklanabilirlik dengesi nedeniyle uygulama tarafında ana model olarak konumlandırılmıştır. Modelin değişken önem düzeyleri üretmesi, SHAP/LIME yorumlarıyla birlikte sunulduğunda karar mekanizmasını daha anlaşılır hale getirir."
    ], 13.8, GRAY)
    n += 1

    slide = content_slide(prs, "Deney Kurulumu ve Ölçüm Mantığı", n)
    add_paragraphs(slide, 0.80, 1.25, 6.18, 3.95, [
        "Veri %75 eğitim ve %25 test olacak şekilde ayrılmıştır. Eğitim kümesi modelin öğrenmesi, test kümesi ise daha önce görülmeyen örneklerde modelin genelleme başarısını ölçmek için kullanılmıştır.",
        "Sınıf dengesizliği nedeniyle tek başına doğruluk yanıltıcı olabilir. Bu nedenle kesinlik, duyarlılık, F1 skoru ve ROC-AUC birlikte değerlendirilmiştir.",
        "Model karşılaştırmasının daha güçlü olması için Random Forest ve çok katmanlı algılayıcı aynı test örnekleri üzerinde McNemar testiyle karşılaştırılmıştır."
    ], 13.6)
    add_table(slide, 7.35, 1.42, 4.85, 2.95, [
        ["Metrik", "Yorum"],
        ["Doğruluk", "Toplam doğru tahmin oranı"],
        ["Kesinlik", "Pozitif tahminlerin ne kadarının doğru olduğu"],
        ["Duyarlılık", "Gerçek pozitiflerin ne kadarının yakalandığı"],
        ["F1", "Kesinlik ve duyarlılık dengesi"],
        ["ROC-AUC", "Sınıfları ayırt etme gücü"],
    ], 9.7)
    n += 1

    slide = content_slide(prs, "Karar Eşiği ve Bankacılık Yorumu", n)
    add_paragraphs(slide, 0.80, 1.25, 6.20, 4.05, [
        "Model çıktısı yalnızca sınıf etiketi değildir; çoğu model müşterinin abone olma olasılığını da üretir. Bu olasılık, karar eşiğine göre düşük veya yüksek ihtimal olarak yorumlanır.",
        "Bankacılıkta karar eşiği iş stratejisine göre değişebilir. Eğer kampanya ekibi daha fazla potansiyel müşteriyi yakalamak istiyorsa duyarlılığı artıran bir eşik tercih edilebilir. Eğer gereksiz aramaları azaltmak daha önemliyse kesinlik odaklı bir eşik kullanılabilir.",
        "Bu nedenle proje, modelin yalnızca teknik doğruluğunu değil, karar destek sistemi olarak nasıl kullanılabileceğini de tartışır."
    ], 13.6)
    add_callout(slide, 7.45, 1.50, 4.60, 1.55, "Örnek kullanım", "Günlük 1.000 müşterilik listede model, en yüksek olasılıklı ilk 200 müşteriyi öncelikli arama listesine alabilir.", BLUE, 13.6, 11.8)
    add_callout(slide, 7.45, 3.45, 4.60, 1.70, "Yorumlama ihtiyacı", "Pazarlama uzmanı, sadece kimin aranacağını değil, modelin bu öneriyi hangi değişkenler nedeniyle yaptığını da görmelidir.", ORANGE, 13.6, 11.8)
    n += 1

    section_slide(prs, "5. Performans Sonuçları", "Sonuçlar, yüksek görünen skorların gerçekçi kullanım koşullarıyla birlikte değerlendirilmesi gerektiğini göstermektedir.", [
        "Görüşme süresi dahil edildiğinde modeller daha yüksek performans gösterebilir; ancak bu bilgi kampanya öncesi bilinmediğinden pratikte kullanılması doğru değildir.",
        "Bu nedenle asıl yorum, görüşme süresi hariç senaryoda elde edilen sonuçlara dayanır."
    ], n, PURPLE)
    n += 1

    slide = content_slide(prs, "Performans Tablosu ve Ana Yorum", n)
    add_table(slide, 0.35, 1.12, 12.65, 3.58, [
        ["Senaryo", "Model", "Doğ.", "Kes.", "Duy.", "F1", "AUC"],
        ["Görüşme süresi dahil", "Random Forest", "0.868", "0.457", "0.925", "0.611", "0.949"],
        ["Görüşme süresi dahil", "Çok Katmanlı Algılayıcı", "0.915", "0.650", "0.528", "0.582", "0.947"],
        ["Görüşme süresi hariç", "Lojistik Regresyon", "0.835", "0.367", "0.644", "0.467", "0.805"],
        ["Görüşme süresi hariç", "Random Forest", "0.864", "0.430", "0.633", "0.512", "0.816"],
        ["Görüşme süresi hariç", "Çok Katmanlı Algılayıcı", "0.900", "0.647", "0.245", "0.355", "0.808"],
    ], 9.8)
    add_paragraphs(slide, 0.82, 5.08, 11.7, 0.95, [
        "Tablo, görüşme süresi dahil senaryoda performansın belirgin şekilde yükseldiğini gösterir. Ancak bu artış gerçek kampanya öncesi kullanım için yanıltıcıdır.",
        "Görüşme süresi hariç senaryoda Random Forest, duyarlılık ve F1 dengesi açısından daha uygulanabilir bir çözüm sunar; bu nedenle ana model olarak tercih edilmiştir."
    ], 12.8, GRAY)
    n += 1

    slide = content_slide(prs, "ROC-AUC Karşılaştırması", n)
    add_image(slide, "model_auc_comparison_turkce.png", 0.90, 1.13, w=6.2)
    add_paragraphs(slide, 7.25, 1.30, 5.05, 3.95, [
        "ROC-AUC metriği, modelin abone olan ve olmayan müşterileri ayırt etme gücünü ölçer. Değer 1’e yaklaştıkça modelin sınıfları ayırma becerisi artar.",
        "Grafikte Random Forest ve çok katmanlı algılayıcının benzer ayırt etme gücüne sahip olduğu görülür. Ancak AUC tek başına model seçimi için yeterli değildir; duyarlılık, kesinlik ve açıklanabilirlik de birlikte değerlendirilmelidir.",
        "Random Forest, uygulama senaryosunda daha dengeli yorumlanabilirlik sağladığı için öne çıkar."
    ], 13.2)
    n += 1

    slide = content_slide(prs, "Hata Matrisi: Random Forest", n)
    add_image(slide, "rf_confusion_matrix.png", 0.95, 1.18, w=4.95)
    add_paragraphs(slide, 6.45, 1.30, 5.85, 3.90, [
        "Hata matrisi, modelin doğru ve yanlış tahminlerini sınıf bazında gösterir. Bu görsel, yalnızca toplam başarıyı değil, hataların hangi yönde yoğunlaştığını anlamaya yardımcı olur.",
        "Bankacılık senaryosunda iki hata türünün maliyeti farklıdır. Abone olabilecek bir müşteriyi düşük ihtimal olarak görmek fırsat kaybına yol açabilir. Abone olmayacak bir müşteriyi yüksek ihtimal olarak görmek ise gereksiz arama maliyeti oluşturur.",
        "Bu nedenle hata matrisi, modelin kampanya stratejisine nasıl entegre edileceğini tartışmak için önemlidir."
    ], 13.2)
    n += 1

    slide = content_slide(prs, "McNemar Testi ile İstatistiksel Değerlendirme", n)
    add_callout(slide, 0.75, 1.25, 5.95, 1.75, "Test sonucu", "Random Forest ve Çok Katmanlı Algılayıcı aynı test örnekleri üzerinde karşılaştırılmıştır.\nb = 450 · c = 819 · χ² = 106.717 · p < 0.001", BLUE, 13.6, 11.9)
    add_paragraphs(slide, 7.10, 1.25, 5.20, 2.78, [
        "McNemar testi, iki sınıflandırıcının aynı örneklerde yaptığı hataların rastlantısal olarak aynı kabul edilip edilemeyeceğini inceler.",
        "p-değerinin 0.001’den küçük olması, iki modelin hata davranışları arasında istatistiksel olarak anlamlı bir fark bulunduğunu gösterir.",
        "Bu sonuç, model karşılaştırmasının yalnızca metrik tablosuna dayanmadığını; istatistiksel analizle desteklendiğini gösterir."
    ], 12.9)
    add_bullets(slide, 0.95, 4.45, 11.4, 1.15, [
        "Yönergedeki bonus beklentilere uygun olarak istatistiksel karşılaştırma yapılmıştır.",
        "Analiz, model seçiminin daha savunulabilir hale gelmesine katkı sağlar."
    ], 13.8)
    n += 1

    slide = content_slide(prs, "Model Seçimi: Neden Random Forest?", n)
    add_table(slide, 0.65, 1.25, 12.0, 2.85, [
        ["Ölçüt", "Random Forest değerlendirmesi"],
        ["Performans", "Görüşme süresi hariç senaryoda dengeli F1 ve duyarlılık sunar."],
        ["Açıklanabilirlik", "Değişken önemleri, SHAP/LIME yorumlarıyla desteklenebilir."],
        ["Kararlılık", "Birden fazla karar ağacını birleştirdiği için tek ağaca göre daha dayanıklıdır."],
        ["Uygulama uyumu", "Streamlit arayüzünde hızlı tahmin üretmeye uygundur."],
    ], 10.9)
    add_paragraphs(slide, 0.85, 4.55, 11.6, 1.02, [
        "Çok katmanlı algılayıcı bazı metriklerde güçlü görünse de pozitif sınıfı yakalama bakımından daha sınırlı kalmıştır. Banka kampanyasında potansiyel müşterileri kaçırmak önemli olduğundan Random Forest daha dengeli ve açıklanabilir bir ana model olarak seçilmiştir."
    ], 13.5, GRAY)
    n += 1

    section_slide(prs, "6. Açıklanabilirlik Analizi", "Modelin ne tahmin ettiği kadar, bu tahmini hangi gerekçelerle yaptığı da incelenmiştir.", [
        "Açıklanabilirlik bölümü, seçilen proje başlığının merkezidir. Kara kutu bir modelin kararlarını kullanıcıya anlaşılır biçimde sunmak, modelin gerçek karar destek sürecinde kullanılabilmesi için gereklidir.",
        "Projede hem model geneli için değişken etkileri hem de tek müşteri örneği için yerel açıklama üretilmiştir."
    ], n, NAVY)
    n += 1

    slide = content_slide(prs, "Küresel Açıklama: Genel Karar Faktörleri", n)
    add_image(slide, "rf_feature_importance_from_explain.png", 0.70, 1.12, w=6.45)
    add_paragraphs(slide, 7.45, 1.20, 4.85, 4.15, [
        "Küresel açıklama, modelin tüm veri seti genelinde hangi değişkenlere daha fazla ağırlık verdiğini gösterir. Bu, modelin genel karar stratejisini anlamak için kullanılır.",
        "Örneğin önceki kampanya sonucu, ekonomik göstergeler veya kampanya temas sayısı gibi değişkenlerin öne çıkması, modelin yalnızca müşteri profili bilgilerine değil, müşteri geçmişi ve dönemsel ekonomik koşullara da baktığını gösterir.",
        "Bu analiz, banka yöneticisinin modelin hangi sinyallere duyarlı olduğunu görmesine yardımcı olur."
    ], 12.9)
    n += 1

    slide = content_slide(prs, "Yerel Açıklama: Tek Müşteri Kararı", n)
    add_image(slide, "lime_explanation_from_explain.png", 0.70, 1.10, w=6.45)
    add_paragraphs(slide, 7.45, 1.18, 4.85, 4.22, [
        "Yerel açıklama, tek bir müşteri için modelin tahminini hangi faktörlerin artırdığını veya azalttığını gösterir. Bu yaklaşım, bireysel müşteri kararlarının açıklanması açısından önemlidir.",
        "Örneğin bir müşterinin önceki kampanyaya olumlu tepki vermesi tahmini artırabilirken, bazı ekonomik göstergeler veya yoğun kampanya teması tahmini düşürebilir.",
        "Böylece kullanıcı yalnızca “yüksek ihtimal” veya “düşük ihtimal” sonucunu görmez; bu sonucun hangi koşullarla oluştuğunu da yorumlayabilir."
    ], 12.9)
    n += 1

    slide = content_slide(prs, "SHAP ve LIME Yaklaşımlarının Rolü", n)
    add_table(slide, 0.65, 1.25, 12.0, 2.65, [
        ["Yöntem", "Projedeki işlevi", "Sağladığı katkı"],
        ["SHAP", "Değişkenlerin tahmine katkısını tutarlı biçimde yorumlama", "Genel ve yerel etki yönlerini açıklamaya yardımcı olur."],
        ["LIME", "Tek bir tahmin çevresinde yerel açıklama üretme", "Belirli müşteri kararının anlaşılmasını kolaylaştırır."],
        ["Değişken önemleri", "Random Forest modelinin genel önceliklerini görme", "Modelin hangi sinyalleri güçlü kullandığını gösterir."],
    ], 10.8)
    add_paragraphs(slide, 0.85, 4.35, 11.65, 1.15, [
        "Bu yöntemler modelin iç yapısını tamamen basitleştirmez; fakat karmaşık karar mekanizmasını insanın yorumlayabileceği sinyallere dönüştürür.",
        "Özellikle finansal karar destek sistemlerinde bu dönüşüm, model çıktısının denetlenebilir ve açıklanabilir olmasını sağlar."
    ], 13.4)
    n += 1

    slide = content_slide(prs, "Açıklanabilirliğin Gerçek Kullanım Değeri", n)
    add_paragraphs(slide, 0.78, 1.25, 6.20, 4.10, [
        "Bir pazarlama uzmanı, modelin yüksek olasılık verdiği müşterileri gördüğünde bu sonucun neden oluştuğunu bilmek ister. Eğer model yalnızca skor üretirse kullanıcı, önerinin iş bilgisiyle uyumlu olup olmadığını kontrol edemez.",
        "XAI çıktıları, modelin önerisini tartışılabilir hale getirir. Örneğin müşteri önceki kampanyada olumlu tepki verdiyse veya belirli ekonomik dönemde daha yüksek eğilim gösteriyorsa bu bilgi karar sürecine ek bağlam sağlar.",
        "Bu nedenle açıklanabilirlik, model güvenini artırır, hatalı veya beklenmedik kararları fark etmeyi kolaylaştırır ve sunum sırasında projenin akademik değerini güçlendirir."
    ], 13.4)
    add_callout(slide, 7.45, 1.45, 4.65, 1.55, "Kullanıcı güveni", "Model çıktısının gerekçeli sunulması, kullanıcıların sistemi daha bilinçli değerlendirmesini sağlar.", TEAL, 13.7, 11.8)
    add_callout(slide, 7.45, 3.40, 4.65, 1.55, "Denetlenebilirlik", "Karar faktörleri görünür olduğunda modelin beklenmedik davranışları daha kolay incelenebilir.", ORANGE, 13.7, 11.8)
    n += 1

    section_slide(prs, "7. Uygulama Arayüzü ve Teslim Çıktıları", "Proje, raporda anlatılan yöntemi çalışır bir Türkçe tahmin sistemine dönüştürür.", [
        "Streamlit arayüzü, hocanın modelin nasıl kullanılacağını doğrudan görmesini sağlar. Kullanıcı müşteri bilgilerini Türkçe alan adlarıyla girer, model tahmin üretir ve açıklanabilirlik çıktıları okunabilir biçimde sunulur.",
        "Bu bölüm, rapor dışındaki uygulamalı teslim değerini gösterir."
    ], n, TEAL)
    n += 1

    slide = content_slide(prs, "Türkçe Tahmin Arayüzü", n)
    add_paragraphs(slide, 0.78, 1.25, 6.20, 3.55, [
        "Arayüz tamamen Türkçe hazırlanmıştır. Meslek, medeni durum, eğitim durumu, konut kredisi, ihtiyaç kredisi, iletişim kanalı ve ekonomik göstergeler gibi alanlar kullanıcıya Türkçe gösterilir.",
        "Tahmin sonucu, teknik sınıf etiketi yerine anlaşılır bir ifade olarak sunulur: “Vadeli mevduata abone olma ihtimali düşük” veya “Vadeli mevduata abone olma ihtimali yüksek”.",
        "Bu tercih, projenin yalnızca kod çalıştıran bir deney değil, kullanıcıya dönük karar destek sistemi olarak görünmesini sağlar."
    ], 13.5)
    add_callout(slide, 7.45, 1.45, 4.65, 1.50, "Çalıştırma", "python train.py\npython explain.py\npython -m streamlit run app.py", BLUE, 13.7, 12.2)
    add_callout(slide, 7.45, 3.35, 4.65, 1.70, "Sunum avantajı", "Hoca, raporda anlatılan modelin gerçek arayüzde nasıl tahmin verdiğini ve açıklama ürettiğini uygulamalı olarak görebilir.", ORANGE, 13.7, 11.8)
    n += 1

    slide = content_slide(prs, "Uygulama Akışı", n)
    add_step_bar(slide, ["Müşteri\nbilgisi", "Ön\nişleme", "Model\ntahmini", "Sonuç\nmetni", "XAI\nçıktısı", "Karar\ndesteği"])
    add_paragraphs(slide, 0.90, 3.75, 11.45, 1.65, [
        "Kullanıcı arayüzde müşteri bilgilerini girer. Bu bilgiler eğitim sırasında kullanılan ön işleme adımlarından geçirilir ve Random Forest modeline aktarılır.",
        "Model olasılık üretir; uygulama bu olasılığı Türkçe bir sonuç metnine dönüştürür. Alt bölümde ise genel karar faktörleri gösterilerek tahminin hangi değişkenlerle ilişkili olduğu açıklanır."
    ], 14.0)
    n += 1

    slide = content_slide(prs, "Üç Kişilik Rol Dağılımı", n)
    add_table(slide, 0.75, 1.25, 11.85, 2.45, [
        ["Öğrenci", "Numara", "Sorumluluk alanı"],
        ["Sude ANLAŞ", "231041038", "Veri hazırlığı, arayüz Türkçeleştirme ve rapor düzeni"],
        ["Doğan Fatih OĞUR", "231041048", "Modelleme, metrik karşılaştırması ve deney kurulumu"],
        ["Merve Naz ORAN", "231041058", "XAI analizi, istatistiksel değerlendirme ve sunum içeriği"],
    ], 11.0)
    add_paragraphs(slide, 0.85, 4.15, 11.65, 1.25, [
        "Rol dağılımı, yalnızca isim tablosu olarak bırakılmamıştır; proje kapsamı üç kişilik ekip çalışmasına uygun biçimde genişletilmiştir.",
        "Veri hazırlığı, modelleme, açıklanabilirlik, istatistiksel test, rapor ve çalışan uygulama çıktıları farklı iş paketleri olarak ele alınmıştır."
    ], 13.8)
    n += 1

    section_slide(prs, "8. Sonuç ve Değerlendirme", "Proje, tahmin performansı ile açıklanabilirliği birlikte ele alan uygulanabilir bir veri madenciliği çalışmasıdır.", [
        "Çalışmanın ana çıktısı, banka kampanyalarında müşteri önceliklendirmesi yapabilen ve model kararlarını Türkçe olarak açıklayabilen bir sistemdir.",
        "Elde edilen sonuçlar, model başarısının tek başına yeterli olmadığını; gerçek kullanım senaryosu, veri sızıntısı, sınıf dengesizliği ve açıklanabilirlik boyutlarının birlikte değerlendirilmesi gerektiğini göstermektedir."
    ], n, NAVY)
    n += 1

    slide = content_slide(prs, "Genel Sonuçlar", n)
    add_bullets(slide, 0.85, 1.25, 11.75, 4.55, [
        "Banka pazarlama problemi, denetimli ikili sınıflandırma problemi olarak modellenmiş ve CRISP-DM sürecine uygun biçimde yürütülmüştür.",
        "Görüşme süresi değişkeninin veri sızıntısı riski taşıdığı gösterilmiş; gerçekçi uygulama değerlendirmesi bu değişken çıkarılarak yapılmıştır.",
        "Random Forest, performans, pozitif sınıfı yakalama ve açıklanabilirlik dengesi nedeniyle ana model olarak seçilmiştir.",
        "SHAP/LIME yaklaşımı ve değişken önemleri, kara kutu model kararlarını anlaşılır hale getirmiştir.",
        "Streamlit arayüzü sayesinde proje, rapor dışında çalışır ve gösterilebilir bir final proje çıktısına dönüştürülmüştür."
    ], 14.0)
    n += 1

    slide = content_slide(prs, "Kısıtlar ve Gelecek Çalışmalar", n)
    add_table(slide, 0.65, 1.25, 12.0, 3.05, [
        ["Kısıt / geliştirme alanı", "Açıklama"],
        ["Veri seti bağlamı", "Veri belirli bir kampanya dönemine ait olduğu için farklı dönemlerde yeniden doğrulama gerekir."],
        ["Sınıf dengesizliği", "Pozitif sınıf azınlıkta olduğu için eşik optimizasyonu ve maliyet duyarlı öğrenme denenebilir."],
        ["Açıklanabilirlik sınırı", "SHAP/LIME model davranışını açıklar; doğrudan nedensellik kanıtı sunmaz."],
        ["Gerçek uygulama", "Canlı banka verisiyle zaman bazlı validasyon ve kullanıcı geri bildirimi eklenebilir."],
    ], 10.6)
    add_paragraphs(slide, 0.85, 4.75, 11.6, 0.92, [
        "Gelecek çalışmalarda karar eşiği kampanya maliyetine göre ayarlanabilir, farklı müşteri segmentleri için ayrı modeller denenebilir ve açıklanabilirlik çıktıları kullanıcı geri bildirimiyle iyileştirilebilir."
    ], 13.5, GRAY)
    n += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, 0.85, 2.15, 11.7, 0.80, "Teşekkürler", 42, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, 1.10, 3.25, 11.1, 0.55, "Sorularınızı alabiliriz.", 22, RGBColor(219, 234, 254), False, PP_ALIGN.CENTER)
    add_paragraphs(slide, 2.00, 4.18, 9.35, 0.82, [
        "Proje; veri hazırlığı, modelleme, istatistiksel değerlendirme, açıklanabilirlik ve Türkçe uygulama arayüzü bileşenleriyle birlikte sunulmuştur."
    ], 14.2, WHITE)
    add_footer(slide, n)
    n += 1

    slide = content_slide(prs, "Kaynaklar", n)
    add_bullets(slide, 0.75, 1.20, 12.0, 4.95, [
        "Moro, S., Cortez, P. ve Rita, P. (2014). A data-driven approach to predict the success of bank telemarketing.",
        "UCI Machine Learning Repository. Bank Marketing Dataset.",
        "Lundberg, S. M. ve Lee, S. I. (2017). A unified approach to interpreting model predictions.",
        "Ribeiro, M. T., Singh, S. ve Guestrin, C. (2016). Why should I trust you? Explaining the predictions of any classifier.",
        "scikit-learn, Streamlit, pandas, matplotlib/seaborn dokümantasyonları."
    ], 14.2)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build()
