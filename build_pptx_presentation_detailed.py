from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "teslim"
OUT.mkdir(exist_ok=True)

PPTX_PATH = OUT / "sunum_3_kisilik_xai.pptx"

NAVY = RGBColor(31, 78, 121)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(15, 118, 110)
ORANGE = RGBColor(249, 115, 22)
GRAY = RGBColor(75, 85, 99)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(17, 24, 39)
RED = RGBColor(220, 38, 38)


def set_text(frame, text, size=22, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_textbox(slide, x, y, w, h, text, size=22, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape.text_frame, text, size, color, bold, align)
    return shape


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.24, 12.2, 0.55, title, 24, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.90), Inches(12.2), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.58, 0.98, 11.9, 0.42, subtitle, 12.5, GRAY)


def add_footer(slide, n):
    add_textbox(slide, 0.55, 7.18, 8.8, 0.22, "BLM308 Veri Madenciliği Final Projesi", 8.5, GRAY)
    add_textbox(slide, 12.2, 7.18, 0.5, 0.22, str(n), 8.5, GRAY, align=PP_ALIGN.RIGHT)


def content_slide(prs, title, n, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    add_footer(slide, n)
    return slide


def section_slide(prs, title, subtitle, detail, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, 0.75, 2.0, 11.9, 0.8, title, 34, WHITE, True)
    add_textbox(slide, 0.8, 3.02, 11.3, 0.55, subtitle, 19, RGBColor(219, 234, 254))
    add_textbox(slide, 0.82, 3.85, 10.9, 0.75, detail, 15, WHITE)
    add_footer(slide, n)
    return slide


def add_bullets(slide, x, y, w, h, bullets, size=17, color=BLACK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return shape


def add_numbered(slide, x, y, w, h, items, size=16):
    return add_bullets(slide, x, y, w, h, [f"{i}. {item}" for i, item in enumerate(items, 1)], size)


def add_callout(slide, x, y, w, h, title, body, color=BLUE, title_size=14, body_size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = color
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(body_size)
    p2.font.color.rgb = BLACK
    p2.space_before = Pt(4)
    return shape


def add_table(slide, x, y, w, h, data, font_size=11):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for r_idx, row in enumerate(data):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r_idx == 0 else RGBColor(248, 250, 252)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE if r_idx == 0 else BLACK
                p.alignment = PP_ALIGN.CENTER
    return shape


def add_image(slide, image_name, x, y, w=None, h=None):
    path = FIG / image_name
    if path.exists():
        kwargs = {}
        if w:
            kwargs["width"] = Inches(w)
        if h:
            kwargs["height"] = Inches(h)
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)
    return add_callout(slide, x, y, w or 5, h or 2, "Görsel bulunamadı", image_name, RED)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    n = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(239, 246, 255)
    add_textbox(slide, 0.72, 0.65, 11.8, 1.45, "Banka Pazarlama Kampanyalarında\nMüşteri Tahmin ve Açıklanabilirlik Sistemi", 33, NAVY, True)
    add_textbox(slide, 0.8, 2.35, 9.7, 0.42, "BLM308 Veri Madenciliği - Final Proje", 18, GRAY)
    add_textbox(slide, 0.8, 3.1, 10.8, 0.75, "Bu sunumda, banka kampanyalarında müşterinin vadeli mevduata abone olma ihtimalini tahmin eden ve model kararlarını anlaşılır hale getiren uçtan uca bir veri madenciliği çalışması anlatılmaktadır.", 15, BLACK)
    add_textbox(slide, 0.8, 5.95, 10.8, 0.34, "Sude ANLAŞ · Doğan Fatih OĞUR · Merve Naz ORAN", 17, BLACK, True)
    add_textbox(slide, 0.8, 6.35, 8.2, 0.28, "Gedik Üniversitesi · Bilgisayar Mühendisliği · Bahar 2026", 12.5, GRAY)
    for i, (label, col) in enumerate([("Veri", BLUE), ("Model", TEAL), ("Açıklama", ORANGE)]):
        x = 8.25 + i * 1.45
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(4.08), Inches(1.05), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        add_textbox(slide, x - 0.05, 5.25, 1.18, 0.28, label, 11.5, NAVY, True, PP_ALIGN.CENTER)
    add_footer(slide, n)
    n += 1

    slide = content_slide(prs, "İçerik", n, "Sunum, rapordaki CRISP-DM akışını izleyerek problemi iş bağlamından teknik sonuçlara doğru adım adım açıklar.")
    add_numbered(slide, 0.9, 1.55, 11.6, 4.95, [
        "Problem ve iş değeri: Bankanın neden böyle bir tahmin sistemine ihtiyaç duyduğu açıklanır.",
        "Veri seti ve hazırlık adımları: Veri kaynağı, sınıf dengesizliği ve veri sızıntısı riski anlatılır.",
        "Modelleme yaklaşımı: Temel modeller ve kara kutu modellerin neden seçildiği özetlenir.",
        "Performans sonuçları: Metrikler, ROC-AUC grafiği, hata matrisi ve McNemar testi yorumlanır.",
        "Açıklanabilirlik çıktıları: Modelin genel ve tek müşteri düzeyindeki karar gerekçeleri gösterilir.",
        "Uygulama ve sonuç: Türkçe arayüz, takım rolleri, kısıtlar ve gelecek çalışma önerileri sunulur.",
    ], 15.2)
    n += 1

    section_slide(prs, "1. Problem", "Amaç yalnızca doğru tahmin yapmak değil, bankanın bu tahmine neden güvenebileceğini de göstermektir.", "Bu bölümde kampanya önceliklendirme probleminin neden önemli olduğu, modelin hangi iş kararını desteklediği ve açıklanabilirliğin finans alanında neden kritik olduğu açıklanır.", n)
    n += 1

    slide = content_slide(prs, "Problem Tanımı", n)
    add_callout(slide, 0.8, 1.35, 11.6, 1.35, "Temel problem", "Banka, pazarlama kampanyasında hangi müşterilerin vadeli mevduata abone olma ihtimalinin yüksek olduğunu önceden tahmin etmek istemektedir. Böylece kampanya ekibi tüm müşterileri aynı öncelikle aramak yerine, dönüşüm potansiyeli daha yüksek müşterilere odaklanabilir.", BLUE, 15, 13.5)
    add_bullets(slide, 1.0, 3.05, 11.0, 2.65, [
        "Rastgele arama yaklaşımı hem zaman kaybına hem de düşük dönüşüm oranına neden olabilir.",
        "Sınırlı çağrı merkezi kapasitesi, veriye dayalı müşteri önceliklendirme ile daha verimli kullanılabilir.",
        "Bankacılıkta kararın açıklanabilir olması önemlidir; bu nedenle modelin yalnızca sonuç değil, gerekçe de üretmesi beklenir.",
    ], 17.5)
    n += 1

    slide = content_slide(prs, "Motivasyon ve İş Değeri", n)
    add_callout(slide, 0.75, 1.35, 5.85, 2.45, "İş değeri", "Model, müşterileri abone olma olasılığına göre sıralayarak kampanya ekibinin daha doğru hedefleme yapmasına yardımcı olur. Bu yaklaşım, arama başına dönüşüm oranını artırabilir ve gereksiz temasları azaltabilir.", TEAL, 15, 13.2)
    add_callout(slide, 6.85, 1.35, 5.75, 2.45, "Açıklanabilirlik ihtiyacı", "Kara kutu bir model yüksek skor üretebilir; ancak banka çalışanı bu skorun hangi müşteri özelliklerinden etkilendiğini görmezse modele güvenmekte zorlanır. XAI çıktıları bu güven boşluğunu azaltır.", ORANGE, 15, 13.2)
    add_bullets(slide, 1.0, 4.4, 11.0, 1.35, [
        "Bu nedenle proje, tahmin başarısını iş değeri ve açıklanabilirlik ile birlikte değerlendiren bir karar destek sistemi olarak tasarlanmıştır.",
    ], 17.5)
    n += 1

    slide = content_slide(prs, "Projenin Güçlü Tarafı", n)
    add_bullets(slide, 0.95, 1.35, 11.4, 4.95, [
        "Çalışma yalnızca bir model eğitme denemesi değildir; veri hazırlığı, modelleme, değerlendirme, açıklanabilirlik ve uygulama arayüzü birlikte ele alınmıştır.",
        "Random Forest ve Çok Katmanlı Algılayıcı gibi kara kutu modeller karşılaştırılmış, model davranışını anlamak için açıklanabilirlik çıktıları üretilmiştir.",
        "Görüşme süresi değişkeninin kampanya öncesinde bilinmediği fark edilerek veri sızıntısı riski ayrıca tartışılmıştır.",
        "Üç kişilik proje kapsamını güçlendirmek için McNemar testi ve Türkçe tahmin arayüzü eklenmiştir.",
    ], 17.2)
    n += 1

    section_slide(prs, "2. Yöntem", "Proje, CRISP-DM metodolojisine uygun biçimde iş probleminden uygulama arayüzüne kadar uçtan uca yürütülmüştür.", "Bu yapı sayesinde çalışma, yalnızca teknik metrikler sunan bir deney olmaktan çıkıp, gerçek bir veri madenciliği projesi gibi iş değeri, veri kalitesi, model güvenilirliği ve sunulabilirlik boyutlarını birlikte kapsar.", n)
    n += 1

    slide = content_slide(prs, "Genel Akış: CRISP-DM", n)
    flow = ["İş\nAnlayışı", "Veri\nAnlayışı", "Veri\nHazırlığı", "Model\nEğitimi", "Değer-\nlendirme", "XAI ve\nArayüz"]
    colors = [NAVY, BLUE, TEAL, ORANGE, RGBColor(124, 58, 237), RGBColor(14, 116, 144)]
    for i, label in enumerate(flow):
        x = 0.55 + i * 2.08
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.40), Inches(1.65), Inches(0.92))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()
        set_text(shape.text_frame, label, 14, WHITE, True, PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_textbox(slide, x + 1.67, 2.67, 0.35, 0.2, "→", 18, GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.95, 4.0, 11.5, 1.05, "CRISP-DM akışı projede yol haritası olarak kullanılmıştır. Önce iş problemi netleştirilmiş, ardından veri incelenmiş, modelleme yapılmış, sonuçlar değerlendirilmiş ve son adımda model kararlarını açıklayan Türkçe bir arayüz hazırlanmıştır.", 17, BLACK, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Kullanılan Araçlar", n)
    add_callout(slide, 0.75, 1.35, 5.85, 2.65, "Modelleme ve analiz", "Python ortamında pandas ve numpy ile veri işleme yapılmış, scikit-learn ile sınıflandırma modelleri kurulmuş, matplotlib ve seaborn ile grafik çıktıları hazırlanmıştır.", BLUE, 15, 13.2)
    add_callout(slide, 6.85, 1.35, 5.75, 2.65, "Açıklanabilirlik ve sunum", "Model açıklamaları SHAP/LIME yaklaşımıyla yorumlanmış, uygulama arayüzü Streamlit ile hazırlanmış, rapor Word formatında ve sunum PowerPoint formatında teslim edilebilir hale getirilmiştir.", TEAL, 15, 13.2)
    add_textbox(slide, 0.95, 4.55, 11.4, 0.72, "Kod yapısı VS Code üzerinde çalıştırılabilecek şekilde düzenlenmiştir. Bu sayede hoca, rapor dışında projenin gerçek çalışan kısmını da görebilir.", 17, BLACK, align=PP_ALIGN.CENTER)
    n += 1

    section_slide(prs, "3. Veri", "Veri seti, hedef değişken dengesizliği ve görüşme süresi kaynaklı veri sızıntısı riski nedeniyle dikkatli hazırlanmıştır.", "Bu bölümde veri kaynağı, hedef değişken, sınıf dağılımı ve ön işleme kararları açıklanır. Özellikle görüşme süresi değişkeni, gerçek kampanya öncesi kullanımda bilinmediği için ayrıca ele alınmıştır.", n)
    n += 1

    slide = content_slide(prs, "Veri Kaynağı ve Tanım", n)
    add_table(slide, 0.8, 1.32, 6.15, 2.25, [
        ["Özellik", "Değer"],
        ["Kaynak", "UCI Bank Marketing"],
        ["Örnek sayısı", "41.188"],
        ["Girdi değişkeni", "20"],
        ["Hedef", "Vadeli mevduat aboneliği"],
    ], 12.5)
    add_callout(slide, 7.25, 1.42, 5.0, 2.15, "Sınıf dağılımı", "Abone oldu: 4.640 (%11,3)\nAbone olmadı: 36.548 (%88,7)\nBu dağılım, abone olan sınıfın azınlıkta olduğunu gösterir.", ORANGE, 14.5, 12.8)
    add_textbox(slide, 0.95, 4.35, 11.6, 0.92, "Veri setindeki hedef değişken, müşterinin kampanya sonunda vadeli mevduata abone olup olmadığını gösterir. Sınıf dengesizliği nedeniyle doğruluk tek başına yanıltıcı olabileceğinden diğer metrikler de raporlanmıştır.", 16.5, BLACK, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Hedef Değişken Dağılımı", n)
    add_image(slide, "target_distribution_turkce.png", 2.15, 1.2, w=9.0)
    add_textbox(slide, 1.0, 6.15, 11.4, 0.55, "Grafik, abone olmayan müşterilerin veri setinde açık biçimde çoğunlukta olduğunu gösterir. Bu nedenle modelin pozitif sınıfı, yani abone olabilecek müşterileri yakalama başarısı özellikle önemlidir.", 15.3, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Veri Hazırlığı ve Kritik Karar", n)
    add_numbered(slide, 0.85, 1.35, 6.4, 3.55, [
        "Kategorik değişkenler makine öğrenmesi modellerinin kullanabileceği sayısal forma dönüştürüldü.",
        "Sayısal değişkenler özellikle sinir ağı modelinin sağlıklı öğrenebilmesi için standartlaştırıldı.",
        "Veri, hedef sınıf dağılımı korunarak eğitim ve test kümelerine ayrıldı.",
        "Sınıf dengesizliğinin etkisini azaltmak için uygun modellerde sınıf ağırlığı kullanıldı.",
    ], 14.6)
    add_callout(slide, 7.55, 1.6, 4.7, 2.75, "Veri sızıntısı farkındalığı", "Görüşme süresi, müşteriyle yapılan arama bittikten sonra bilinir. Kampanya öncesinde bu bilgi elde olmadığı için gerçekçi uygulama modelinde kullanılmamıştır.", RED, 14.5, 12.8)
    n += 1

    section_slide(prs, "4. Modelleme", "Temel modeller ve kara kutu modeller birlikte kullanılarak hem karşılaştırma hem de açıklanabilirlik zemini oluşturulmuştur.", "Model seçimi, ders kapsamındaki sınıflandırma mantığını ve proje başlığındaki XAI gereksinimini aynı anda karşılayacak şekilde yapılmıştır.", n)
    n += 1

    slide = content_slide(prs, "Seçilen Modeller", n)
    add_table(slide, 0.65, 1.3, 12.05, 3.25, [
        ["Model", "Tür", "Projedeki rolü"],
        ["Lojistik Regresyon", "Temel model", "Basit ve yorumlanabilir referans"],
        ["Karar Ağacı", "Ağaç tabanlı", "Kural benzeri karşılaştırma"],
        ["Random Forest", "Ensemble kara kutu", "Ana tahmin ve açıklama modeli"],
        ["Çok Katmanlı Algılayıcı", "Sinir ağı", "Alternatif kara kutu model"],
    ], 12.2)
    add_textbox(slide, 0.9, 5.25, 11.6, 0.65, "Random Forest, hem gerçekçi senaryoda dengeli sonuç verdiği hem de değişken önem düzeyleriyle açıklanabildiği için uygulamada ana model olarak seçilmiştir.", 16.5, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Deney Kurulumu", n)
    add_bullets(slide, 1.0, 1.35, 11.2, 4.85, [
        "Veri %75 eğitim ve %25 test olacak şekilde ayrılmıştır; böylece model daha önce görmediği örneklerde değerlendirilmiştir.",
        "Stratified split kullanılarak abone olan ve olmayan sınıfların oranı eğitim ve test kümelerinde korunmuştur.",
        "Doğruluk, kesinlik, duyarlılık, F1 ve ROC-AUC birlikte raporlanmıştır; çünkü tek bir metrik bu problem için yeterli değildir.",
        "Random Forest ve Çok Katmanlı Algılayıcı arasındaki fark, aynı test örnekleri üzerinde McNemar testiyle ayrıca incelenmiştir.",
    ], 16.7)
    n += 1

    section_slide(prs, "5. Sonuçlar", "Görüşme süresi hariç gerçekçi senaryoda Random Forest dengeli ve açıklanabilir ana model olarak öne çıkmıştır.", "Bu bölümde metrik tablosu, ROC-AUC grafiği, hata matrisi ve McNemar testi birlikte yorumlanır. Amaç, sadece en yüksek skoru değil, iş açısından en anlamlı modeli seçmektir.", n)
    n += 1

    slide = content_slide(prs, "Performans Karşılaştırması", n)
    add_table(slide, 0.35, 1.18, 12.6, 3.60, [
        ["Senaryo", "Model", "Doğ.", "Kes.", "Duy.", "F1", "AUC"],
        ["Görüşme süresi dahil", "Random Forest", "0.868", "0.457", "0.925", "0.611", "0.949"],
        ["Görüşme süresi dahil", "Çok Katmanlı Algılayıcı", "0.915", "0.650", "0.528", "0.582", "0.947"],
        ["Görüşme süresi hariç", "Lojistik Regresyon", "0.835", "0.367", "0.644", "0.467", "0.805"],
        ["Görüşme süresi hariç", "Random Forest", "0.864", "0.430", "0.633", "0.512", "0.816"],
        ["Görüşme süresi hariç", "Çok Katmanlı Algılayıcı", "0.900", "0.647", "0.245", "0.355", "0.808"],
    ], 10.2)
    add_textbox(slide, 0.85, 5.25, 11.7, 0.72, "Görüşme süresi dahil senaryo daha yüksek skor üretir; fakat bu değişken kampanya öncesinde bilinmediği için gerçekçi kullanımda görüşme süresi hariç sonuçlar esas alınmıştır.", 15.5, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "ROC-AUC Karşılaştırması", n)
    add_image(slide, "model_auc_comparison_turkce.png", 1.75, 1.16, w=9.8)
    add_textbox(slide, 0.95, 6.18, 11.6, 0.56, "ROC-AUC, modelin iki sınıfı ayırt etme gücünü gösterir. Random Forest, gerçekçi senaryoda güçlü bir ayrım performansı sunarken açıklanabilirlik açısından da avantaj sağlamıştır.", 15.2, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Hata Matrisi: Random Forest", n)
    add_image(slide, "rf_confusion_matrix.png", 4.15, 1.23, w=4.85)
    add_textbox(slide, 1.0, 6.15, 11.4, 0.62, "Hata matrisi, modelin hangi sınıflarda daha başarılı olduğunu gösterir. Banka açısından abone olabilecek müşterileri kaçırmak önemli bir kayıp olabileceği için duyarlılık değeri ayrıca izlenmelidir.", 15.1, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "İstatistiksel Karşılaştırma", n)
    add_callout(slide, 0.95, 1.35, 11.4, 1.35, "McNemar testi: Random Forest vs Çok Katmanlı Algılayıcı", "b = 450     c = 819     χ² = 106.717     p < 0.001\nBu sonuç, iki modelin aynı test örnekleri üzerindeki hata örüntülerinin anlamlı biçimde farklılaştığını gösterir.", BLUE, 14.5, 12.7)
    add_bullets(slide, 1.05, 3.15, 11.1, 2.55, [
        "Bu test, yalnızca metrik değerlerine bakmak yerine iki modelin hangi örneklerde doğru veya yanlış karar verdiğini karşılaştırır.",
        "p-değerinin çok düşük olması, Random Forest ve Çok Katmanlı Algılayıcı modellerinin hata davranışlarının rastlantısal olarak aynı kabul edilemeyeceğini gösterir.",
        "Bu nedenle değerlendirme bölümü, yönergede beklenen istatistiksel analiz boyutunu da karşılar.",
    ], 15.8)
    n += 1

    section_slide(prs, "6. Açıklanabilirlik", "Kara kutu model kararları, hem genel faktör etkileri hem de tek müşteri düzeyindeki açıklamalarla yorumlanmıştır.", "Bu bölüm, projenin XAI başlığıyla doğrudan ilişkili kısmıdır. Modelin hangi değişkenlerden etkilendiği ve belirli bir müşteriye neden yüksek ya da düşük olasılık verdiği açıklanır.", n)
    n += 1

    slide = content_slide(prs, "Modelin Genel Kararlarını Etkileyen Faktörler", n)
    add_image(slide, "rf_feature_importance_from_explain.png", 0.95, 1.08, w=11.45)
    add_textbox(slide, 0.85, 6.18, 11.7, 0.62, "Grafik, Random Forest modelinin genel olarak hangi değişkenlere daha fazla ağırlık verdiğini gösterir. Ekonomik göstergeler, önceki kampanya sonucu ve kampanya temas bilgileri karar mekanizmasında öne çıkmaktadır.", 14.8, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Tek Müşteri İçin Yerel Açıklama", n)
    add_image(slide, "lime_explanation_from_explain.png", 0.95, 1.08, w=11.45)
    add_textbox(slide, 0.85, 6.18, 11.7, 0.62, "Yerel açıklama, tek bir müşteri için model kararını etkileyen faktörleri gösterir. Böylece modelin yalnızca sonuç üretmesi değil, bu sonuca hangi koşullar nedeniyle yaklaştığı da anlaşılır.", 14.8, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Açıklanabilirlik Neden Önemli?", n)
    add_bullets(slide, 1.0, 1.35, 11.2, 4.95, [
        "Kara kutu modeller karmaşık ilişkileri öğrenebilir; ancak karar gerekçesi görünmediğinde kullanıcı güveni azalabilir.",
        "Açıklanabilirlik çıktıları, pazarlama uzmanının model önerisini iş bilgisiyle birlikte değerlendirmesine imkân tanır.",
        "Finans alanında kararların şeffaf, denetlenebilir ve gerekçelendirilebilir olması beklenir.",
        "Bu nedenle XAI, projede ek bir süs değil, modelin gerçek hayatta kullanılabilirliğini artıran temel bileşendir.",
    ], 17)
    n += 1

    section_slide(prs, "7. Uygulama ve Sonuç", "Türkçe arayüz, proje çıktısını hocaya gösterilebilecek çalışan bir karar destek sistemine dönüştürür.", "Bu son bölümde arayüzün nasıl çalıştığı, takım rollerinin nasıl dağıtıldığı, projenin genel sonucu ve gelecek çalışma önerileri özetlenir.", n)
    n += 1

    slide = content_slide(prs, "Türkçe Tahmin Arayüzü", n)
    add_numbered(slide, 0.85, 1.32, 6.15, 3.55, [
        "Kullanıcı müşteri bilgilerini Türkçe alan adlarıyla girer.",
        "Model müşterinin vadeli mevduata abone olma olasılığını hesaplar.",
        "Sonuç, düşük veya yüksek ihtimal ifadesiyle anlaşılır biçimde gösterilir.",
        "Alt bölümde modelin genel kararlarını etkileyen faktörler listelenir.",
    ], 15.2)
    add_callout(slide, 7.35, 1.45, 4.85, 2.65, "Sunum dili", "Arayüzde teknik olmayan İngilizce ifadeler kullanılmaz. Müşteri alanları, sonuç metinleri ve açıklanabilirlik tablosu Türkçe hazırlanmıştır.", TEAL, 14.5, 12.8)
    add_textbox(slide, 0.95, 5.45, 11.5, 0.48, "Çalıştırma komutu: python -m streamlit run app.py", 17, NAVY, True, PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Takım İçi Rol Dağılımı", n)
    add_table(slide, 0.85, 1.45, 11.6, 2.35, [
        ["Öğrenci", "Numara", "Rol"],
        ["Sude ANLAŞ", "231041038", "Veri mühendisi"],
        ["Doğan Fatih OĞUR", "231041048", "Model mühendisi"],
        ["Merve Naz ORAN", "231041058", "Değerlendirme ve XAI analisti"],
    ], 13)
    add_textbox(slide, 0.9, 4.45, 11.6, 0.82, "Üç kişilik ekip kapsamı, yalnızca kişi sayısı artırılarak değil; veri hazırlığı, modelleme, istatistiksel değerlendirme, XAI yorumu ve uygulama arayüzü olarak iş bölümü genişletilerek karşılanmıştır.", 16.2, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Sonuç", n)
    add_bullets(slide, 1.0, 1.3, 11.2, 5.0, [
        "Banka pazarlama problemi CRISP-DM sürecine uygun olarak problem tanımından uygulama arayüzüne kadar uçtan uca ele alınmıştır.",
        "Görüşme süresi değişkeni için veri sızıntısı farkındalığı gösterilmiş ve gerçekçi senaryoda bu değişken modelden çıkarılmıştır.",
        "Random Forest, gerçekçi senaryoda performans ve açıklanabilirlik dengesi nedeniyle ana model olarak seçilmiştir.",
        "McNemar testi, model karşılaştırmasının istatistiksel açıdan da desteklenmesini sağlamıştır.",
        "Türkçe arayüz sayesinde proje rapor dışında uygulamalı olarak da gösterilebilir hale getirilmiştir.",
    ], 16.5)
    n += 1

    slide = content_slide(prs, "Kısıtlar ve Gelecek Çalışmalar", n)
    add_callout(slide, 0.75, 1.3, 5.85, 2.95, "Kısıtlar", "Veri seti belirli bir ülke ve kampanya dönemine aittir. Abone olan müşteri sınıfı azınlıktadır. Açıklanabilirlik çıktıları nedensellik göstermez; yalnızca model davranışını yorumlamaya yardımcı olur.", ORANGE, 14.5, 12.8)
    add_callout(slide, 6.85, 1.3, 5.75, 2.95, "Gelecek çalışmalar", "İleride maliyet duyarlı öğrenme, karar eşiği optimizasyonu, zaman bazlı validasyon ve farklı açıklanabilirlik yöntemlerinin karşılaştırılması yapılabilir.", TEAL, 14.5, 12.8)
    add_textbox(slide, 0.95, 5.15, 11.4, 0.55, "Bu geliştirmeler, modelin gerçek banka ortamına daha güvenilir biçimde taşınmasına katkı sağlayabilir.", 16, GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, 0.9, 2.35, 11.5, 0.8, "Teşekkürler", 42, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, 3.42, 11.5, 0.5, "Sorularınız?", 24, RGBColor(219, 234, 254), align=PP_ALIGN.CENTER)
    add_footer(slide, n)
    n += 1

    slide = content_slide(prs, "Kaynaklar", n)
    add_bullets(slide, 0.75, 1.2, 12.0, 4.95, [
        "Moro, Cortez ve Rita (2014), banka telepazarlama kampanyalarının veri madenciliğiyle tahmini.",
        "UCI Machine Learning Repository, Bank Marketing Dataset.",
        "Lundberg ve Lee (2017), SHAP yaklaşımı.",
        "Ribeiro, Singh ve Guestrin (2016), LIME yaklaşımı.",
    ], 17)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build()
