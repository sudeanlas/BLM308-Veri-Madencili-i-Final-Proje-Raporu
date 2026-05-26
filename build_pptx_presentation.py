from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
FIG = ROOT / "figures"
OUT = ROOT / "teslim"
OUT.mkdir(exist_ok=True)

PPTX_PATH = OUT / "sunum_3_kisilik_xai.pptx"

NAVY = RGBColor(31, 78, 121)
BLUE = RGBColor(37, 99, 235)
TEAL = RGBColor(15, 118, 110)
ORANGE = RGBColor(249, 115, 22)
GRAY = RGBColor(75, 85, 99)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(17, 24, 39)
RED = RGBColor(220, 38, 38)


def set_text(frame, text, size=24, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_textbox(slide, x, y, w, h, text, size=24, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box.text_frame, text, size=size, color=color, bold=bold, align=align)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.55, 0.28, 12.2, 0.55, title, size=25, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.93), Inches(12.2), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.58, 1.02, 11.9, 0.32, subtitle, size=12.5, color=GRAY)


def add_footer(slide, n):
    add_textbox(slide, 0.55, 7.18, 8.5, 0.22, "BLM308 Veri Madenciliği Final Projesi", size=8.5, color=GRAY)
    add_textbox(slide, 12.2, 7.18, 0.5, 0.22, str(n), size=8.5, color=GRAY, align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, bullets, size=19, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_numbered(slide, x, y, w, h, items, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items, start=1):
        p = tf.paragraphs[0] if i == 1 else tf.add_paragraph()
        p.text = f"{i}. {item}"
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(7)
    return box


def add_callout(slide, x, y, w, h, title, body, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = color
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Aptos"
    p2.font.size = Pt(14)
    p2.font.color.rgb = BLACK
    return shape


def add_image(slide, image_name, x, y, w=None, h=None):
    path = FIG / image_name
    if path.exists():
        kwargs = {}
        if w is not None:
            kwargs["width"] = Inches(w)
        if h is not None:
            kwargs["height"] = Inches(h)
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)
    add_callout(slide, x, y, w or 5, h or 3, "Görsel bulunamadı", image_name, RED)
    return None


def add_table(slide, x, y, w, h, data, font_size=12):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r_idx, row in enumerate(data):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE if r_idx == 0 else BLACK
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r_idx == 0 else RGBColor(248, 250, 252)
    return table_shape


def section_slide(prs, title, subtitle, n):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    add_textbox(slide, 0.72, 2.3, 11.8, 0.75, title, size=34, color=WHITE, bold=True)
    add_textbox(slide, 0.78, 3.25, 10.8, 0.65, subtitle, size=19, color=RGBColor(219, 234, 254))
    add_footer(slide, n)
    return slide


def content_slide(prs, title, n, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    add_footer(slide, n)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    n = 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(239, 246, 255)
    add_textbox(slide, 0.75, 0.65, 11.6, 1.45, "Banka Pazarlama Kampanyalarında\nMüşteri Tahmin ve Açıklanabilirlik Sistemi", size=34, color=NAVY, bold=True)
    add_textbox(slide, 0.8, 2.35, 8.5, 0.38, "BLM308 Veri Madenciliği - Final Proje", size=18, color=GRAY)
    add_textbox(slide, 0.8, 5.95, 10.8, 0.34, "Sude ANLAŞ · Doğan Fatih OĞUR · Merve Naz ORAN", size=17, color=BLACK, bold=True)
    add_textbox(slide, 0.8, 6.35, 8.2, 0.28, "Gedik Üniversitesi · Bilgisayar Mühendisliği · Bahar 2026", size=12.5, color=GRAY)
    # Distinct visual motif: a simple decision path.
    for i, (label, col) in enumerate([("Veri", BLUE), ("Model", TEAL), ("Açıklama", ORANGE)]):
        x = 8.25 + i * 1.45
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.55), Inches(1.05), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = col
        shape.line.fill.background()
        add_textbox(slide, x - 0.05, 4.75, 1.18, 0.28, label, size=11.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, n)
    n += 1

    slide = content_slide(prs, "İçerik", n)
    add_numbered(slide, 1.1, 1.45, 10.8, 4.8, [
        "Problem ve iş değeri",
        "Veri seti ve hazırlık adımları",
        "Modelleme yaklaşımı",
        "Performans sonuçları ve istatistiksel test",
        "Açıklanabilirlik çıktıları",
        "Türkçe tahmin arayüzü ve sonuç",
    ], size=20)
    n += 1

    section_slide(prs, "1. Problem", "Amaç yalnızca tahmin yapmak değil, karar gerekçesini de göstermek.", n)
    n += 1

    slide = content_slide(prs, "Problem Tanımı", n)
    add_callout(slide, 0.8, 1.45, 11.6, 1.05, "Temel problem", "Banka, pazarlama kampanyasında hangi müşterilerin vadeli mevduata abone olma ihtimalinin yüksek olduğunu önceden tahmin etmek istemektedir.", BLUE)
    add_bullets(slide, 1.05, 2.9, 10.9, 2.8, [
        "Rastgele arama stratejisi zaman ve maliyet kaybı oluşturabilir.",
        "Sınırlı çağrı merkezi kapasitesi daha verimli kullanılmalıdır.",
        "Model yalnızca tahmin üretmemeli; karar gerekçesini de açıklamalıdır.",
    ], size=20)
    n += 1

    slide = content_slide(prs, "Motivasyon ve İş Değeri", n)
    add_callout(slide, 0.8, 1.4, 5.65, 2.0, "İş değeri", "Yüksek potansiyelli müşterilere öncelik verilir, kampanya dönüşüm oranı artırılır ve gereksiz aramalar azaltılır.", TEAL)
    add_callout(slide, 6.85, 1.4, 5.65, 2.0, "Açıklanabilirlik ihtiyacı", "Bankacılıkta model önerisinin gerekçesi denetlenebilir olmalı ve pazarlama ekibi tarafından yorumlanabilmelidir.", ORANGE)
    add_bullets(slide, 1.05, 4.05, 10.7, 1.55, [
        "Tahmin sistemi çağrı önceliklendirme aracı olarak konumlandırılmıştır.",
        "Açıklanabilirlik katmanı modeli karar destek sistemine dönüştürür.",
    ], size=19)
    n += 1

    slide = content_slide(prs, "Proje Farkı", n)
    add_bullets(slide, 0.95, 1.35, 11.4, 4.9, [
        "Uçtan uca CRISP-DM süreci izlenmiştir.",
        "Random Forest ve Çok Katmanlı Algılayıcı kara kutu model olarak karşılaştırılmıştır.",
        "Görüşme süresi değişkeni için veri sızıntısı analizi yapılmıştır.",
        "Random Forest ve Çok Katmanlı Algılayıcı McNemar testiyle karşılaştırılmıştır.",
        "Tamamen Türkçe tahmin ve açıklanabilirlik arayüzü hazırlanmıştır.",
    ], size=20)
    n += 1

    section_slide(prs, "2. Yöntem", "CRISP-DM akışıyla veri, model, değerlendirme ve XAI birlikte ele alındı.", n)
    n += 1

    slide = content_slide(prs, "Genel Akış: CRISP-DM", n)
    flow = ["İş\nAnlayışı", "Veri\nAnlayışı", "Veri\nHazırlığı", "Model\nEğitimi", "Değer-\nlendirme", "XAI ve\nArayüz"]
    colors = [NAVY, BLUE, TEAL, ORANGE, RGBColor(124, 58, 237), RGBColor(14, 116, 144)]
    for i, label in enumerate(flow):
        x = 0.55 + i * 2.08
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.65), Inches(1.65), Inches(0.92))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()
        set_text(shape.text_frame, label, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_textbox(slide, x + 1.67, 2.92, 0.35, 0.2, "→", size=18, color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.45, 4.35, 10.4, 0.6, "Model çıktıları yalnızca metrik olarak değil, iş değeri ve açıklanabilirlik açısından da yorumlandı.", size=18, color=BLACK, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Kullanılan Araçlar", n)
    add_callout(slide, 0.85, 1.45, 5.5, 2.25, "Modelleme ve analiz", "Python, pandas, numpy, scikit-learn, matplotlib ve seaborn kullanıldı.", BLUE)
    add_callout(slide, 6.9, 1.45, 5.5, 2.25, "Açıklanabilirlik ve sunum", "SHAP/LIME yaklaşımı, Streamlit arayüzü, Beamer ve Word raporu hazırlandı.", TEAL)
    add_bullets(slide, 1.0, 4.25, 10.9, 1.3, [
        "Kod yapısı VS Code üzerinde çalıştırılabilir şekilde düzenlenmiştir.",
        "Sunum ve rapor, yönergede beklenen teslim yapısını destekler.",
    ], size=18)
    n += 1

    section_slide(prs, "3. Veri", "UCI Bank Marketing veri seti sınıf dengesizliği ve veri sızıntısı açısından dikkatli ele alındı.", n)
    n += 1

    slide = content_slide(prs, "Veri Kaynağı ve Tanım", n)
    add_table(slide, 0.9, 1.4, 6.0, 2.2, [
        ["Özellik", "Değer"],
        ["Kaynak", "UCI Bank Marketing"],
        ["Örnek sayısı", "41.188"],
        ["Girdi değişkeni", "20"],
        ["Hedef", "Vadeli mevduat aboneliği"],
    ], font_size=13)
    add_callout(slide, 7.25, 1.55, 4.95, 1.75, "Sınıf dağılımı", "Abone oldu: 4.640 (%11,3)\nAbone olmadı: 36.548 (%88,7)", ORANGE)
    add_textbox(slide, 7.25, 4.0, 4.8, 0.95, "Bu dengesizlik nedeniyle doğruluk tek başına yeterli değildir.", size=18, color=BLACK)
    n += 1

    slide = content_slide(prs, "Hedef Değişken Dağılımı", n)
    add_image(slide, "target_distribution_turkce.png", 2.15, 1.35, w=9.0)
    add_textbox(slide, 1.45, 6.35, 10.5, 0.35, "Abone olan sınıf azınlıkta olduğu için duyarlılık ve kesinlik birlikte değerlendirilmiştir.", size=16, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Veri Hazırlığı ve Kritik Karar", n)
    add_numbered(slide, 0.9, 1.35, 6.3, 3.6, [
        "Kategorik değişkenler sayısallaştırıldı.",
        "Sayısal değişkenler standartlaştırıldı.",
        "Veri eğitim/test olarak ayrıldı.",
        "Sınıf dengesizliği için ağırlık kullanıldı.",
    ], size=18)
    add_callout(slide, 7.45, 1.75, 4.8, 2.35, "Veri sızıntısı farkındalığı", "Görüşme süresi müşteri aranmadan önce bilinmez. Bu nedenle gerçekçi kullanımda modelden çıkarılmıştır.", RED)
    n += 1

    section_slide(prs, "4. Modelleme", "Temel modeller ve kara kutu modeller birlikte karşılaştırıldı.", n)
    n += 1

    slide = content_slide(prs, "Seçilen Modeller", n)
    add_table(slide, 0.7, 1.35, 11.9, 3.2, [
        ["Model", "Tür", "Projedeki rolü"],
        ["Lojistik Regresyon", "Temel model", "Yorumlanabilir karşılaştırma"],
        ["Karar Ağacı", "Ağaç tabanlı", "Kural benzeri referans"],
        ["Random Forest", "Ensemble kara kutu", "Ana tahmin modeli"],
        ["Çok Katmanlı Algılayıcı", "Sinir ağı", "Alternatif kara kutu model"],
    ], font_size=12.5)
    add_textbox(slide, 1.0, 5.35, 11.3, 0.45, "Random Forest, performans ve açıklanabilirlik dengesi nedeniyle arayüzde kullanılan ana modeldir.", size=17, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Deney Kurulumu", n)
    add_bullets(slide, 1.05, 1.45, 10.8, 4.4, [
        "Veri bölme: %75 eğitim, %25 test",
        "Bölme yöntemi: Stratified split",
        "Ana metrikler: doğruluk, kesinlik, duyarlılık, F1, ROC-AUC",
        "Ek değerlendirme: McNemar testi",
        "Uygulama senaryosu: görüşme süresi hariç Random Forest",
    ], size=20)
    n += 1

    section_slide(prs, "5. Sonuçlar", "Gerçekçi senaryoda Random Forest dengeli ve açıklanabilir ana model olarak öne çıktı.", n)
    n += 1

    slide = content_slide(prs, "Performans Karşılaştırması", n)
    add_table(slide, 0.35, 1.25, 12.6, 3.55, [
        ["Senaryo", "Model", "Doğ.", "Kes.", "Duy.", "F1", "AUC"],
        ["Görüşme süresi dahil", "Random Forest", "0.868", "0.457", "0.925", "0.611", "0.949"],
        ["Görüşme süresi dahil", "Çok Katmanlı Algılayıcı", "0.915", "0.650", "0.528", "0.582", "0.947"],
        ["Görüşme süresi hariç", "Lojistik Regresyon", "0.835", "0.367", "0.644", "0.467", "0.805"],
        ["Görüşme süresi hariç", "Random Forest", "0.864", "0.430", "0.633", "0.512", "0.816"],
        ["Görüşme süresi hariç", "Çok Katmanlı Algılayıcı", "0.900", "0.647", "0.245", "0.355", "0.808"],
    ], font_size=10.2)
    add_textbox(slide, 1.0, 5.45, 11.2, 0.42, "Gerçekçi kullanımda görüşme süresi hariç sonuçlar esas alınmıştır.", size=16, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "ROC-AUC Karşılaştırması", n)
    add_image(slide, "model_auc_comparison_turkce.png", 1.75, 1.25, w=9.8)
    add_textbox(slide, 1.2, 6.25, 10.9, 0.42, "Görüşme süresi dahil senaryo performansı yükseltir; gerçek kampanya öncesi kullanım için yanıltıcı olabilir.", size=15.5, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Hata Matrisi: Random Forest", n)
    add_image(slide, "rf_confusion_matrix.png", 4.25, 1.32, w=4.7)
    add_textbox(slide, 1.35, 6.22, 10.7, 0.45, "Model abone olmayan sınıfta güçlüdür; abone olabilecek müşterileri kaçırmamak için duyarlılık ayrıca izlenmelidir.", size=15.5, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "İstatistiksel Karşılaştırma", n)
    add_callout(slide, 1.0, 1.55, 11.2, 1.15, "McNemar testi: Random Forest vs Çok Katmanlı Algılayıcı", "b = 450     c = 819     χ² = 106.717     p < 0.001", BLUE)
    add_bullets(slide, 1.1, 3.25, 10.8, 2.2, [
        "İki model aynı test örnekleri üzerinde karşılaştırıldı.",
        "p-değeri çok düşük olduğu için hata örüntüleri istatistiksel olarak farklıdır.",
        "Değerlendirme yalnızca metrik tablosuna bağlı bırakılmadı.",
    ], size=19)
    n += 1

    section_slide(prs, "6. Açıklanabilirlik", "Kara kutu model kararları genel ve yerel düzeyde yorumlandı.", n)
    n += 1

    slide = content_slide(prs, "Modelin Genel Kararlarını Etkileyen Faktörler", n)
    add_image(slide, "rf_feature_importance_from_explain.png", 0.95, 1.18, w=11.45)
    add_textbox(slide, 1.1, 6.32, 11.0, 0.35, "Ekonomik göstergeler, önceki kampanya sonucu ve kampanya temas bilgileri öne çıkmaktadır.", size=15.5, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Tek Müşteri İçin Yerel Açıklama", n)
    add_image(slide, "lime_explanation_from_explain.png", 0.95, 1.15, w=11.45)
    add_textbox(slide, 1.1, 6.3, 11.0, 0.35, "Yerel açıklama, hangi faktörlerin abone olma ihtimalini artırıp azalttığını gösterir.", size=15.5, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Açıklanabilirlik Neden Önemli?", n)
    add_bullets(slide, 1.05, 1.45, 10.8, 4.2, [
        "Kara kutu model kararları denetlenebilir hale gelir.",
        "Pazarlama uzmanı model önerisini iş bağlamında yorumlayabilir.",
        "Finans alanında güven, şeffaflık ve hesap verebilirlik artar.",
        "Model yalnızca skor üreten bir araç değil, karar destek sistemi olur.",
    ], size=20)
    n += 1

    section_slide(prs, "7. Uygulama ve Sonuç", "Türkçe arayüz proje çıktısını uygulamalı hale getiriyor.", n)
    n += 1

    slide = content_slide(prs, "Türkçe Tahmin Arayüzü", n)
    add_numbered(slide, 0.95, 1.35, 5.8, 3.1, [
        "Müşteri bilgileri girilir.",
        "Model abone olma olasılığını üretir.",
        "Karar düşük/yüksek ihtimal olarak gösterilir.",
        "Genel etkili faktörler listelenir.",
    ], size=17)
    add_callout(slide, 7.1, 1.45, 5.0, 2.15, "Sunum dili", "Tüm alanlar Türkçedir. Teknik olmayan İngilizce ifadeler kaldırılmıştır. Görüşme süresi ana modelde kullanılmaz.", TEAL)
    add_textbox(slide, 1.05, 5.45, 11.1, 0.42, "Çalıştırma: python -m streamlit run app.py", size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Takım İçi Rol Dağılımı", n)
    add_table(slide, 1.0, 1.55, 11.3, 2.35, [
        ["Öğrenci", "Numara", "Rol"],
        ["Sude ANLAŞ", "231041038", "Veri mühendisi"],
        ["Doğan Fatih OĞUR", "231041048", "Model mühendisi"],
        ["Merve Naz ORAN", "231041058", "Değerlendirme ve XAI analisti"],
    ], font_size=13)
    add_textbox(slide, 1.25, 4.6, 10.8, 0.55, "Üç kişilik ekip kapsamı; modelleme, uygulama, XAI ve istatistiksel test ile genişletilmiştir.", size=18, color=GRAY, align=PP_ALIGN.CENTER)
    n += 1

    slide = content_slide(prs, "Sonuç", n)
    add_bullets(slide, 1.05, 1.35, 10.8, 4.75, [
        "Banka pazarlama problemi CRISP-DM süreciyle uçtan uca ele alındı.",
        "Görüşme süresi değişkeni için veri sızıntısı farkındalığı gösterildi.",
        "Random Forest gerçekçi senaryoda dengeli ve açıklanabilir ana model olarak seçildi.",
        "McNemar testiyle model farkı istatistiksel olarak desteklendi.",
        "Türkçe arayüzle proje uygulamalı hale getirildi.",
    ], size=19)
    n += 1

    slide = content_slide(prs, "Kısıtlar ve Gelecek Çalışmalar", n)
    add_callout(slide, 0.85, 1.35, 5.7, 2.45, "Kısıtlar", "Veri belirli bir ülke ve kampanya dönemine aittir. Abone olan sınıf azınlıktadır. Açıklamalar nedensellik değil model davranışı gösterir.", ORANGE)
    add_callout(slide, 6.95, 1.35, 5.55, 2.45, "Gelecek çalışmalar", "Maliyet duyarlı öğrenme, karar eşiği optimizasyonu ve farklı XAI yöntemlerinin karşılaştırılması yapılabilir.", TEAL)
    n += 1

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    add_textbox(slide, 0.9, 2.45, 11.5, 0.8, "Teşekkürler", size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.9, 3.45, 11.5, 0.5, "Sorularınız?", size=24, color=RGBColor(219, 234, 254), align=PP_ALIGN.CENTER)
    add_footer(slide, n)
    n += 1

    slide = content_slide(prs, "Kaynaklar", n)
    add_bullets(slide, 0.75, 1.25, 12.0, 4.9, [
        "Moro, Cortez ve Rita (2014), Bank telemarketing başarısının veri madenciliğiyle tahmini.",
        "UCI Machine Learning Repository, Bank Marketing Dataset.",
        "Lundberg ve Lee (2017), SHAP yaklaşımı.",
        "Ribeiro, Singh ve Guestrin (2016), LIME yaklaşımı.",
    ], size=17)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build()
